import logging
import json
import os
from typing import Optional
from docx import Document
import pdfplumber
from pptx import Presentation 
import ffmpeg
from pydub import AudioSegment
import torch
from transformers import WhisperProcessor, WhisperForConditionalGeneration
from fastapi import FastAPI, WebSocket, WebSocketDisconnect
import librosa
import numpy as np
import asyncio
import whisperx
import math
import gc
import pandas as pd
from dotenv import load_dotenv

load_dotenv()


logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO)

AudioSegment.converter = os.getenv("FFMPEG_PATH")
if not AudioSegment.converter:
    logger.error("FFMPEG_PATH environment variable is not set")
    raise ValueError("FFMPEG_PATH is not set")

def extract_text_from_pdf(file_path: str) -> str:
    """
    Извлекает текст из PDF-файла постранично с помощью pdfplumber
    """
    try:
        text = ""
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {e}")
        raise e

def extract_text_from_docx(file_path: str) -> str:
    """
    Извлекает текст из DOCX-файла (Word)
    """
    try:
        doc = Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text)
    except Exception as e:
        logger.error(f"Error extracting text from DOCX: {e}")
        raise e

def extract_text_from_txt(file_path: str) -> str:
    """
    Извлекает текст из текстового файла .txt
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        logger.error(f"Error extracting text from TXT: {e}")
        raise e

def extract_text_from_pptx(file_path: str) -> str:
    """
    Извлекает текст из PPTX (PowerPoint)
    """
    try:
        presentation = Presentation(file_path)
        full_text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return '\n'.join(full_text)
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {e}")
        raise e

class AudioTranscription:
    """
    Класс для транскрипции аудио- и видеофайлов с возможностью разделения на сегменты,
    выравнивания, а также спикер-диаризации (определения говорящих)
    Использует whisperx + ffmpeg + pydub
    """
    def __init__(self, 
                 model_name: str = os.getenv("WHISPER_MODEL_NAME", "large-v3-turbo"), 
                 compute_type: str = "float16", 
                 batch_size: int = 16, 
                 hf_token: str = os.getenv("HF_TOKEN", ""),
                 diarization: bool = False):
        """
        Инициализация класса, загрузка моделей WhisperX и опционально модели диаризации
        """
        self.device = "cuda" if torch.cuda.is_available() else "cpu"
        self.batch_size = batch_size
        self.compute_type = compute_type
        self.hf_token = hf_token
        self.segment_length = 30
        self.diarization = diarization 
        if torch.__version__ < "2.6":
            logger.error("PyTorch version must be at least 2.6 due to security issues with torch.load")
            raise ValueError("PyTorch version must be at least 2.6 due to security issues with torch.load")
        logger.info(f"PyTorch version: {torch.__version__}")
        if not self.hf_token and diarization:
            logger.error("HF_TOKEN environment variable is not set for diarization")
            raise ValueError("HF_TOKEN is required for diarization")

        try:
            self.model = whisperx.load_model(model_name, self.device, compute_type=compute_type, language="ru")
            logger.info(f"WhisperX model loaded: {model_name}")
        except Exception as e:
            logger.error(f"Error loading WhisperX model: {str(e)}")
            raise e

        try:
            self.align_model, self.metadata = whisperx.load_align_model(language_code="ru", device=self.device)
            logger.info("Alignment model loaded")
        except Exception as e:
            logger.error(f"Error loading alignment model: {str(e)}")
            raise e

        if self.diarization:
            try:
                self.diarization_pipeline = whisperx.diarize.DiarizationPipeline(use_auth_token=self.hf_token, device=self.device)
                logger.info("Diarization pipeline loaded")
            except Exception as e:
                logger.error(f"Error loading diarization pipeline: {str(e)}")
                raise e
        else:
            self.diarization_pipeline = None
            logger.info("Diarization is disabled")

    def extract_audio_from_video(self, video_path: str, audio_path: str) -> None:
        """
        Извлекает аудио-дорожку из видеофайла и сохраняет её в формате WAV
        """
        try:
            (
                ffmpeg
                .input(video_path)
                .output(audio_path, format='wav', acodec='pcm_s16le', ac=1, ar='16k')
                .overwrite_output()
                .run(quiet=True)
            )
            logger.info(f"Audio extracted from {video_path} to {audio_path}")
        except Exception as e:
            logger.error(f"Error extracting audio: {str(e)}")
            raise e

    def convert_audio_to_wav(self, input_path: str, output_path: str) -> None:
        """
        Конвертирует любой аудиофайл в WAV 16kHz моно
        """
        try:
            audio = AudioSegment.from_file(input_path)
            audio = audio.set_frame_rate(16000).set_channels(1)
            audio.export(output_path, format="wav")
            logger.info(f"Audio converted to WAV: {output_path}")
        except Exception as e:
            logger.error(f"Error converting to WAV: {str(e)}")
            raise e

    def split_audio(self, audio_file: str) -> tuple[list, int]:
        """
        Разбивает аудио на сегменты длиной 30 секунд для батчевой обработки
        """
        try:
            audio = AudioSegment.from_file(audio_file)
            duration = len(audio) / 1000
            num_segments = math.ceil(duration / self.segment_length)
            segments = []
            for i in range(num_segments):
                start_time = i * self.segment_length * 1000 
                end_time = min((i + 1) * self.segment_length * 1000, len(audio))
                segment = audio[start_time:end_time]
                segment_file = f"segment_{i}.wav"
                segment.export(segment_file, format="wav")
                segments.append((segment_file, i * self.segment_length))
            return segments, num_segments
        except Exception as e:
            logger.error(f"Error splitting audio: {str(e)}")
            return [], 0

    def transcribe(self, audio_path: str, progress_callback=None, diarization: Optional[bool] = None) -> str:
        """
        Транскрибирует аудиофайл с опцией выравнивания и диаризации
        """
        if diarization is None:
            diarization = self.diarization

        try:
            segments, num_segments = self.split_audio(audio_path)
            if not segments:
                logger.error("Failed to split audio into segments")
                return ""

            all_segments = []
            for i, (segment_file, start_offset) in enumerate(segments):
                try:
                    audio_segment = whisperx.load_audio(segment_file)
                    result = self.model.transcribe(audio_segment, batch_size=self.batch_size, language="ru")
                    for seg in result["segments"]:
                        seg["start"] += start_offset
                        seg["end"] += start_offset
                        all_segments.append(seg)
                    os.remove(segment_file)
                    if progress_callback:
                        if diarization:
                            progress_callback((i + 1) / num_segments * 0.7)
                        else:
                            progress_callback((i + 1) / num_segments * 0.8)
                except Exception as e:
                    logger.error(f"Error processing segment {i + 1}: {str(e)}")
                finally:
                    if os.path.exists(segment_file):
                        os.remove(segment_file)

            transcription_result = {"segments": all_segments, "language": "ru"}

            if progress_callback:
                if diarization:
                    progress_callback(0.7)
                else:
                    progress_callback(0.8)

            aligned_result = whisperx.align(transcription_result["segments"], self.align_model, self.metadata, audio_path, self.device, return_char_alignments=False)
            logger.info("Alignment completed")

            if progress_callback:
                if diarization:
                    progress_callback(0.8)
                else:
                    progress_callback(0.9)

            if diarization:
                if self.diarization_pipeline is None:
                    logger.error("Diarization pipeline is not loaded. Set diarization=True during initialization")
                    raise ValueError("Diarization pipeline is not loaded")
                diarization_result = self.diarization_pipeline(audio_path)
                result = whisperx.assign_word_speakers(diarization_result, aligned_result)
                speaker_durations = {}
                for _, row in diarization_result.iterrows():
                    speaker = row['speaker']
                    duration = row['end'] - row['start']
                    speaker_durations[speaker] = speaker_durations.get(speaker, 0) + duration
                dominant_speaker = max(speaker_durations, key=speaker_durations.get)
                dominant_words = [
                    word for segment in result["segments"] for word in segment.get("words", [])
                    if "speaker" in word and word["speaker"] == dominant_speaker and "start" in word and "end" in word
                ]
                dominant_words.sort(key=lambda x: x["start"])
                transcription = " ".join(word["word"] for word in dominant_words if "word" in word)
            else:
                all_words = [word for segment in aligned_result["segments"] for word in segment.get("words", []) if "word" in word]
                all_words.sort(key=lambda x: x.get("start", 0))
                transcription = " ".join(word["word"] for word in all_words)

            if progress_callback:
                progress_callback(1.0)
            return transcription
        except Exception as e:
            logger.error(f"Error transcribing audio: {e}")
            raise e