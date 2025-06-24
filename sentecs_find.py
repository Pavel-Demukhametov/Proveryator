import os
import re
import nltk
from functools import lru_cache
from nltk.corpus import stopwords
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lex_rank import LexRankSummarizer
from transformers import AutoModelForTokenClassification, AutoTokenizer
import torch
from natasha import Segmenter, NewsEmbedding, MorphVocab, NewsMorphTagger, Doc
from pptx import Presentation
from docx import Document
import pdfplumber
import spacy
from nltk.corpus import stopwords

# Инициализация
nltk.download('punkt')
stop_words = set(stopwords.words('russian'))
segmenter = Segmenter()
emb = NewsEmbedding()
morph_vocab = MorphVocab()
morph_tagger = NewsMorphTagger(emb)
nlp = spacy.load('ru_core_news_sm')

# Пути
NER_MODEL_PATH = r'D:\Program Files\Lecture_test_front\fast\ruSciBert-NER'
file_path = r'D:\Аврискин\Java\01. Введение в Java.pptx'

# Загрузка NER модели и токенизатора
tokenizer = AutoTokenizer.from_pretrained(NER_MODEL_PATH, use_fast=True, add_prefix_space=True)
term_model = AutoModelForTokenClassification.from_pretrained(NER_MODEL_PATH)
term_model.eval()

@lru_cache(maxsize=1024)
def lemmatize_text(text: str) -> str:
    doc = Doc(text)
    doc.segment(segmenter)
    doc.tag_morph(morph_tagger)
    for token in doc.tokens:
        token.lemmatize(morph_vocab)
    return ' '.join(token.lemma for token in doc.tokens)

def extract_text(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext == '.pdf':
        with pdfplumber.open(path) as pdf:
            return '\n'.join(page.extract_text() or '' for page in pdf.pages)
    if ext == '.docx':
        doc = Document(path)
        return '\n'.join(p.text for p in doc.paragraphs)
    if ext == '.pptx':
        prs = Presentation(path)
        return '\n'.join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text'))
    raise ValueError("Unsupported file type")

def preprocess(text: str) -> str:
    return re.sub(r'\s+', ' ', text).strip()

def split_into_sentences(text: str, min_w=2, max_w=30) -> list:
    raw = nltk.sent_tokenize(text)
    return [s.strip() for s in raw if min_w <= len(s.split()) <= max_w and any(c.isalpha() for c in s)]

def extract_terms_batch(sentences: list) -> dict:
    inputs = tokenizer(sentences, return_tensors='pt', truncation=True, padding=True, max_length=512)
    inputs = {k: v.to(term_model.device) for k, v in inputs.items()}
    with torch.no_grad():
        outputs = term_model(**inputs).logits
    preds = torch.argmax(outputs, dim=2)
    results = {}
    for idx, sent in enumerate(sentences):
        labels = preds[idx].tolist()
        tokens = tokenizer.convert_ids_to_tokens(inputs['input_ids'][idx])
        terms, curr = [], []
        for tok, lab in zip(tokens, labels):
            tag = term_model.config.id2label[lab]
            if tag == 'B':
                if curr:
                    terms.append(curr)
                curr = [tok]
            elif tag == 'I' and curr:
                curr.append(tok)
            else:
                if curr:
                    terms.append(curr)
                curr = []
        if curr:
            terms.append(curr)
        clean = [tokenizer.decode(tokenizer.convert_tokens_to_ids(t)).replace('▁','').strip().lower() for t in terms]
        lem = lemmatize_text(sent).lower()
        results[sent] = [t for t in set(clean) if t in lem]
    return results

def filter_by_terms(sentences: list, term_map: dict, min_terms=2) -> list:
    return [s for s in sentences if len(term_map.get(s, [])) >= min_terms]

def is_noise(s: str) -> bool:
    doc = nlp(s)
    if not any(tok.pos_ == 'VERB' for tok in doc):
        return True
    words = [tok.text for tok in doc if tok.is_alpha]
    if words and sum(len(w) <= 2 for w in words) / len(words) > 0.3:
        return True
    if re.search(r'(www\.|https?://|\.com|\.ru)', s):
        return True
    if re.search(r'\b[А-ЯЁ]{4,}\b', s):
        return True
    return False

def flesch_score(s: str) -> float:
    words = s.split()
    syllables = sum(len(re.findall(r'[аеёиоуыэюя]', w, re.IGNORECASE)) for w in words)
    n = len(words)
    return 206.835 - 1.015 * n - 84.6 * (syllables / n) if n else 0

def compute_lexrank(sentences: list) -> dict:
    parser = PlaintextParser.from_string(' '.join(sentences), Tokenizer('russian'))
    summarizer = LexRankSummarizer()
    summarizer.stop_words = stop_words  # Установка русских стоп-слов
    summary = summarizer(parser.document, len(sentences))
    n = len(sentences)
    return {str(s): (n - i) / n for i, s in enumerate(summary)}

def compute_novelty_scores(sentences: list) -> dict:
    seen = set()
    nov = {}
    for sent in sentences:
        words = set(sent.split())
        new = words - seen
        nov[sent] = len(new)
        seen |= words
    maxv = max(nov.values()) or 1
    return {s: v / maxv for s, v in nov.items()}

def rank_by_lexrank_then_novelty(sentences: list, lex: dict) -> list:
    novelty = compute_novelty_scores(sentences)

    def sort_key(s):
        return (-round(lex.get(s, 0), 2), -novelty.get(s, 0))

    return sorted(sentences, key=sort_key)

def main(path: str) -> list:
    text = preprocess(extract_text(path))
    sents = split_into_sentences(text)

    term_map = extract_terms_batch(sents)
    candidates = filter_by_terms(sents, term_map, min_terms=1)
    candidates = [s for s in candidates if not is_noise(s)]
    lex_scores = compute_lexrank(candidates)
    ranked = rank_by_lexrank_then_novelty(candidates, lex_scores)
    return ranked

if __name__ == '__main__':
    top = main(file_path)
    print("Отобранные предложения:")
    for idx, s in enumerate(top, 1):
        print(f"{idx}. {s}")
